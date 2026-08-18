from pathlib import Path
from typing import Callable

from pydantic import ValidationError

from skytrap.core.context import WorkspaceContext
from skytrap.tools.base import Tool, ToolResult
from skytrap.tools.filesystem import resolve_in_workspace
from skytrap.tools.registry import RegistryContext, register_tool
from skytrap.tools.skills.pitch_deck.schema import PitchDeckInput, PitchSlide

# Standard python-pptx default template layout indices.
TITLE_SLIDE_LAYOUT = 0
BULLETS_SLIDE_LAYOUT = 1
SECTION_SLIDE_LAYOUT = 2


def _build_preview(parsed: PitchDeckInput) -> str:
    lines = [f"NEW PITCH DECK: {parsed.path}", f"{len(parsed.slides)} slide(s)", ""]
    for i, slide in enumerate(parsed.slides, start=1):
        lines.append(f"{i}. [{slide.type}] {slide.title}")
        if slide.subtitle:
            lines.append(f"   {slide.subtitle}")
        for bullet in slide.bullets:
            lines.append(f"   - {bullet}")
    return "\n".join(lines)


def _add_slide(presentation, slide_spec: PitchSlide) -> None:
    if slide_spec.type == "title":
        slide = presentation.slides.add_slide(presentation.slide_layouts[TITLE_SLIDE_LAYOUT])
        slide.shapes.title.text = slide_spec.title
        if slide_spec.subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = slide_spec.subtitle
    elif slide_spec.type == "section":
        slide = presentation.slides.add_slide(presentation.slide_layouts[SECTION_SLIDE_LAYOUT])
        slide.shapes.title.text = slide_spec.title
    else:  # "bullets"
        slide = presentation.slides.add_slide(presentation.slide_layouts[BULLETS_SLIDE_LAYOUT])
        slide.shapes.title.text = slide_spec.title
        if slide_spec.bullets:
            body = slide.placeholders[1].text_frame
            body.text = slide_spec.bullets[0]
            for bullet in slide_spec.bullets[1:]:
                paragraph = body.add_paragraph()
                paragraph.text = bullet


class PitchDeckTool(Tool):
    name = "pitch_deck"
    description = (
        "Generate a PowerPoint (.pptx) deck from structured slides (title, section, or "
        "bullets). Shows a preview and requires confirmation before writing, same as "
        "write_file. "
        'Arguments: {"path": "<output path ending in .pptx>", "slides": [{"type": '
        '"title"|"section"|"bullets", "title": "...", "subtitle": "...", "bullets": [...]}]}'
    )

    def __init__(self, confirm: Callable[[str], bool]):
        self._confirm = confirm

    def execute(self, workspace: WorkspaceContext, arguments: dict) -> ToolResult:
        try:
            parsed = PitchDeckInput.model_validate(arguments)
        except ValidationError as exc:
            return ToolResult(success=False, output=f"Invalid arguments: {exc}")

        if not parsed.path.lower().endswith(".pptx"):
            return ToolResult(success=False, output="Output path must end in .pptx")

        ok, resolved = resolve_in_workspace(workspace, parsed.path)
        if not ok:
            return ToolResult(success=False, output=resolved)

        preview = _build_preview(parsed)
        if not self._confirm(preview):
            return ToolResult(success=False, output="User declined the write; file not created.")

        from pptx import Presentation

        presentation = Presentation()
        for slide_spec in parsed.slides:
            _add_slide(presentation, slide_spec)

        file_path = Path(resolved)
        file_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(file_path))

        return ToolResult(success=True, output=f"Wrote {len(parsed.slides)} slide(s) to {parsed.path}")


@register_tool
def _build_pitch_deck_tool(context: RegistryContext) -> Tool:
    return PitchDeckTool(confirm=context.confirm_write)
