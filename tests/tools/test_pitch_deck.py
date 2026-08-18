from pathlib import Path

from pptx import Presentation

from skytrap.core.context import WorkspaceContext
from skytrap.tools.skills.pitch_deck.tool import PitchDeckTool


def _workspace(tmp_path: Path) -> WorkspaceContext:
    return WorkspaceContext(path=tmp_path.resolve(), name=tmp_path.name, is_git=False)


DECK_ARGS = {
    "path": "deck.pptx",
    "slides": [
        {"type": "title", "title": "Acme Corp", "subtitle": "Series A pitch"},
        {"type": "section", "title": "The Problem"},
        {
            "type": "bullets",
            "title": "Why now",
            "bullets": ["Market grew 3x in 2 years", "No incumbent solution", "Team has shipped this before"],
        },
    ],
}


def test_confirmed_write_creates_real_pptx_with_correct_content(tmp_path):
    previews = []

    def confirm(preview: str) -> bool:
        previews.append(preview)
        return True

    tool = PitchDeckTool(confirm=confirm)
    result = tool.execute(_workspace(tmp_path), DECK_ARGS)

    assert result.success
    file_path = tmp_path / "deck.pptx"
    assert file_path.exists()

    presentation = Presentation(str(file_path))
    assert len(presentation.slides) == 3

    slide1, slide2, slide3 = presentation.slides
    assert slide1.shapes.title.text == "Acme Corp"
    assert slide1.placeholders[1].text == "Series A pitch"

    assert slide2.shapes.title.text == "The Problem"

    assert slide3.shapes.title.text == "Why now"
    bullet_texts = [p.text for p in slide3.placeholders[1].text_frame.paragraphs]
    assert bullet_texts == [
        "Market grew 3x in 2 years",
        "No incumbent solution",
        "Team has shipped this before",
    ]

    assert "Acme Corp" in previews[0]
    assert "Why now" in previews[0]
    assert "Market grew 3x in 2 years" in previews[0]


def test_declined_write_creates_no_file(tmp_path):
    tool = PitchDeckTool(confirm=lambda preview: False)
    result = tool.execute(_workspace(tmp_path), DECK_ARGS)
    assert not result.success
    assert "declined" in result.output.lower()
    assert not (tmp_path / "deck.pptx").exists()


def test_rejects_non_pptx_extension(tmp_path):
    tool = PitchDeckTool(confirm=lambda preview: True)
    args = {**DECK_ARGS, "path": "deck.txt"}
    result = tool.execute(_workspace(tmp_path), args)
    assert not result.success
    assert ".pptx" in result.output


def test_bullets_slide_with_no_bullets_does_not_crash(tmp_path):
    tool = PitchDeckTool(confirm=lambda preview: True)
    args = {"path": "deck.pptx", "slides": [{"type": "bullets", "title": "Empty for now"}]}
    result = tool.execute(_workspace(tmp_path), args)
    assert result.success

    presentation = Presentation(str(tmp_path / "deck.pptx"))
    (slide,) = presentation.slides
    assert slide.shapes.title.text == "Empty for now"


def test_path_outside_workspace_is_rejected(tmp_path):
    tool = PitchDeckTool(confirm=lambda preview: True)
    args = {**DECK_ARGS, "path": "../../etc/evil.pptx"}
    result = tool.execute(_workspace(tmp_path), args)
    assert not result.success
    assert "outside the workspace" in result.output


def test_invalid_arguments(tmp_path):
    tool = PitchDeckTool(confirm=lambda preview: True)
    result = tool.execute(_workspace(tmp_path), {"path": "deck.pptx", "slides": []})
    assert not result.success
    assert "Invalid arguments" in result.output


def test_creates_parent_directories(tmp_path):
    tool = PitchDeckTool(confirm=lambda preview: True)
    args = {**DECK_ARGS, "path": "nested/dir/deck.pptx"}
    result = tool.execute(_workspace(tmp_path), args)
    assert result.success
    assert (tmp_path / "nested" / "dir" / "deck.pptx").exists()
